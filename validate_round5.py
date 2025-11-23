import sys
import json
import subprocess
from pptx import Presentation
from pptx.util import Inches

def create_multi_master_pptx(filename):
    prs = Presentation()
    # Add a slide to the first master
    slide_layout = prs.slide_layouts[0]
    prs.slides.add_slide(slide_layout)
    
    # Duplicate the master (hacky way since python-pptx doesn't support adding masters easily)
    # Actually, python-pptx doesn't support adding masters directly. 
    # But we can try to load a template that has multiple masters if we had one.
    # Since we don't, we might have to skip the "create" part and just trust the logic fix 
    # if we can't easily reproduce it.
    # HOWEVER, the suggestion says "If samples/multi_master.pptx truly has only one SlideMaster...".
    # I will try to create a file with 2 masters if possible, but it's hard with basic python-pptx.
    # Let's assume for now I can't easily create one, but I can verify the CODE LOGIC change.
    
    prs.save(filename)
    print(f"Created {filename} (single master for now)")

def run_probe(args):
    cmd = ["uv", "run", "tools/ppt_capability_probe.py"] + args
    result = subprocess.run(cmd, capture_output=True, text=True)
    return result

def validate_round5():
    print("Validating Round 5 Suggestions...")
    
    # 1. Check for metadata.masters (should be missing currently)
    result = run_probe(["--file", "Presentation.pptx", "--json"])
    if result.returncode == 0:
        data = json.loads(result.stdout)
        if "masters" in data["metadata"]:
             print("✅ metadata.masters present")
        else:
             print("❌ metadata.masters MISSING")
             
        # Check analysis_complete in capabilities
        if "analysis_complete" in data["capabilities"]:
            print("✅ capabilities.analysis_complete present")
        else:
            print("❌ capabilities.analysis_complete MISSING")

if __name__ == "__main__":
    validate_round5()
