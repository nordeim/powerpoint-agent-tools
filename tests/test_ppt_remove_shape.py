#!/usr/bin/env python3
"""Tests for ppt_remove_shape.py v3.0"""

import pytest
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).parent.parent / "tools"))

from ppt_remove_shape import (
    remove_shape,
    remove_shapes_batch,
    get_shape_details,
    find_shape_by_name,
    validate_removal
)


@pytest.fixture
def sample_pptx_with_shapes(tmp_path):
    """Create a sample PowerPoint with multiple shapes for testing."""
    pptx_path = tmp_path / "test_remove.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add multiple shapes
    shape1 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(1), Inches(2), Inches(1)
    )
    shape1.name = "Rectangle 1"
    
    shape2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(4), Inches(1), Inches(1), Inches(1)
    )
    shape2.name = "Circle 1"
    
    shape3 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(3), Inches(2), Inches(1)
    )
    shape3.name = "Rectangle 2"
    shape3.text_frame.text = "Important content"
    
    prs.save(str(pptx_path))
    return pptx_
