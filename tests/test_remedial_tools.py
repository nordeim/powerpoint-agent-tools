import pytest
import json
import os
import sys
from pathlib import Path
from pptx import Presentation

# Allow importing tools directly
sys.path.insert(0, str(Path(__file__).parent.parent))

from tools.ppt_add_notes import add_notes
from tools.ppt_set_z_order import set_z_order
from tools.ppt_replace_text import replace_text

def test_add_notes(tmp_path):
    """Test adding speaker notes functionality."""
    pptx_path = tmp_path / "test_notes.pptx"
    prs = Presentation()
    # Use a layout that exists (index 0 usually Title Slide)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(pptx_path)
    
    # 1. Add initial note
    res1 = add_notes(pptx_path, 0, "First note", mode="append")
    assert res1["status"] == "success"
    
    # Verify
    prs = Presentation(pptx_path)
    assert prs.slides[0].notes_slide.notes_text_frame.text == "First note"
    
    # 2. Append note
    res2 = add_notes(pptx_path, 0, "Second note", mode="append")
    
    # Verify Append
    prs = Presentation(pptx_path)
    assert "First note" in prs.slides[0].notes_slide.notes_text_frame.text
    assert "Second note" in prs.slides[0].notes_slide.notes_text_frame.text
    
    # 3. Overwrite note
    res3 = add_notes(pptx_path, 0, "Overwritten", mode="overwrite")
    
    # Verify Overwrite
    prs = Presentation(pptx_path)
    assert prs.slides[0].notes_slide.notes_text_frame.text == "Overwritten"

def test_z_order(tmp_path):
    """Test shape layering (Z-Order)."""
    pptx_path = tmp_path / "test_z.pptx"
    prs = Presentation()
    # Use blank layout (usually index 6)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Create shapes (0=Back by default creation order)
    s1 = slide.shapes.add_shape(1, 0, 0, 100, 100)
    s1.name = "Shape 1"
    s2 = slide.shapes.add_shape(1, 10, 10, 100, 100)
    s2.name = "Shape 2" # On top
    
    prs.save(pptx_path)
    
    # Verify initial state (Shape 2 is last in list, so on top)
    prs = Presentation(pptx_path)
    assert prs.slides[0].shapes[-1].name == "Shape 2"
    
    # Action: Send Shape 2 (index 1) to back
    set_z_order(pptx_path, 0, 1, "send_to_back")
    
    # Verify swap
    prs = Presentation(pptx_path)
    # Shape 2 should now be at index 0 (back)
    assert prs.slides[0].shapes[0].name == "Shape 2"
    # Shape 1 should now be at index 1 (front)
    assert prs.slides[0].shapes[1].name == "Shape 1"

def test_targeted_text_replace(tmp_path):
    """Test surgical text replacement."""
    pptx_path = tmp_path / "test_replace.pptx"
    prs = Presentation()
    
    # Slide 0
    slide0 = prs.slides.add_slide(prs.slide_layouts[6])
    tb0 = slide0.shapes.add_textbox(0,0,100,100)
    tb0.text_frame.text = "Target"
    
    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb1 = slide1.shapes.add_textbox(0,0,100,100)
    tb1.text_frame.text = "Target"
    
    prs.save(pptx_path)
    
    # Action: Replace "Target" -> "Hit" ONLY on Slide 1
    replace_text(pptx_path, "Target", "Hit", slide_index=1)
    
    # Verify
    prs = Presentation(pptx_path)
    # Slide 0 should be UNCHANGED
    assert prs.slides[0].shapes[0].text_frame.text == "Target"
    # Slide 1 should be CHANGED
    assert prs.slides[1].shapes[0].text_frame.text == "Hit"
