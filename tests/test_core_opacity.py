import unittest
import os
import sys
from pptx import Presentation
from pptx.util import Inches

# Add core directory to path
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), '../core')))
from powerpoint_agent_core import PowerPointAgent, PowerPointAgentError

class TestOpacity(unittest.TestCase):
    def setUp(self):
        self.test_file = "test_opacity.pptx"
        # Create a blank presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        prs.save(self.test_file)

    def tearDown(self):
        if os.path.exists(self.test_file):
            os.remove(self.test_file)

    def test_add_shape_with_opacity(self):
        with PowerPointAgent(self.test_file) as agent:
            # Add a shape with 50% opacity
            result = agent.add_shape(
                slide_index=0,
                shape_type="rectangle",
                position={"left": 1.0, "top": 1.0},
                size={"width": 2.0, "height": 2.0},
                fill_color="#FF0000",
                fill_opacity=0.5,
                line_color="#0000FF",
                line_opacity=0.5
            )
            
            self.assertTrue(result["styling"]["fill_opacity_applied"])
            self.assertTrue(result["styling"]["line_opacity_applied"])
            self.assertEqual(result["styling"]["fill_opacity"], 0.5)

    def test_format_shape_opacity(self):
        with PowerPointAgent(self.test_file) as agent:
            # Add opaque shape
            add_res = agent.add_shape(
                slide_index=0,
                shape_type="rectangle",
                position={"left": 4.0, "top": 1.0},
                size={"width": 2.0, "height": 2.0},
                fill_color="#00FF00"
            )
            shape_idx = add_res["shape_index"]
            
            # Update to 20% opacity
            fmt_res = agent.format_shape(
                slide_index=0,
                shape_index=shape_idx,
                fill_opacity=0.2
            )
            
            self.assertTrue(fmt_res["changes_detail"]["fill_opacity_applied"])
            self.assertEqual(fmt_res["changes_detail"]["fill_opacity"], 0.2)

if __name__ == '__main__':
    unittest.main()
