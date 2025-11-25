#!/usr/bin/env python3
"""Tests for ppt_format_shape.py v3.0"""

import pytest
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).parent.parent / "tools"))

from ppt_format_shape import (
    format_shape,
    resolve_color,
    resolve_transparency,
    validate_formatting_params,
    COLOR_PRESETS,
    TRANSPARENCY_PRESETS
)


@pytest.fixture
def sample_pptx_with_shapes(tmp_path):
    """Create a sample PowerPoint with shapes for testing."""
    pptx_path = tmp_path / "test_shapes.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add a rectangle
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(1), Inches(2), Inches(1)
    )
    shape.text = "Test Shape"
    
    # Add another shape
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(4), Inches(1), Inches(1), Inches(1)
    )
    
    prs.save(str(pptx_path))
    return pptx_path


class TestResolveColor:
    """Tests for color resolution."""
    
    def test_preset_colors(self):
        assert resolve_color("primary") == "#0070C0"
        assert resolve_color("danger") == "#C00000"
        assert resolve_color("white") == "#FFFFFF"
    
    def test_hex_passthrough(self):
        assert resolve_color("#FF0000") == "#FF0000"
        assert resolve_color("#abc123") == "#abc123"
    
    def test_hex_without_hash(self):
        assert resolve_color("FF0000") == "#FF0000"
    
    def test_transparent(self):
        assert resolve_color("transparent") is None
        assert resolve_color("none") is None
    
    def test_none(self):
        assert resolve_color(None) is None


class TestResolveTransparency:
    """Tests for transparency resolution."""
    
    def test_preset_values(self):
        assert resolve_transparency("opaque") == 0.0
        assert resolve_transparency("subtle") == 0.15
        assert resolve_transparency("medium") == 0.5
    
    def test_numeric_values(self):
        assert resolve_transparency(0.5) == 0.5
        assert resolve_transparency("0.3") == 0.3
    
    def test_percentage(self):
        assert resolve_transparency("30%") == 0.3
        assert resolve_transparency("50%") == 0.5
    
    def test_none(self):
        assert resolve_transparency(None) is None


class TestValidateFormattingParams:
    """Tests for parameter validation."""
    
    def test_valid_params(self):
        result = validate_formatting_params(
            fill_color="#0070C0",
            line_color="#000000",
            transparency=0.3
        )
        assert result["validation_results"]["fill_color_valid"]
        assert len(result["warnings"]) == 0
    
    def test_low_contrast_warning(self):
        result = validate_formatting_params(
            fill_color="#FFFFFF",
            line_color=None,
            transparency=None
        )
        assert len(result["warnings"]) > 0
        assert "contrast" in result["warnings"][0].lower()
    
    def test_invalid_transparency_warning(self):
        result = validate_formatting_params(
            fill_color="#0070C0",
            line_color=None,
            transparency=1.5  # Out of range
        )
        assert any("range" in w.lower() for w in result["warnings"])


class TestFormatShape:
    """Tests for format_shape function."""
    
    def test_format_fill_color(self, sample_pptx_with_shapes):
        result = format_shape(
            filepath=sample_pptx_with_shapes,
            slide_index=0,
            shape_index=0,
            fill_color="#FF0000"
        )
        
        assert result["status"] in ("success", "warning")
        assert result["formatting_applied"]["fill_color"] == "#FF0000"
    
    def test_format_with_transparency(self, sample_pptx_with_shapes):
        result = format_shape(
            filepath=sample_pptx_with_shapes,
            slide_index=0,
            shape_index=0,
            fill_color="#000000",
            transparency=0.5
        )
        
        assert result["status"] in ("success", "warning")
        assert result["formatting_applied"]["transparency"] == 0.5
    
    def test_format_line(self, sample_pptx_with_shapes):
        result = format_shape(
            filepath=sample_pptx_with_shapes,
            slide_index=0,
            shape_index=0,
            line_color="#000000",
            line_width=2.0
        )
        
        assert result["formatting_applied"]["line_color"] == "#000000"
        assert result["formatting_applied"]["line_width"] == 2.0
    
    def test_no_options_raises(self, sample_pptx_with_shapes):
        with pytest.raises(ValueError) as exc_info:
            format_shape(
                filepath=sample_pptx_with_shapes,
                slide_index=0,
                shape_index=0
            )
        assert "at least one" in str(exc_info.value).lower()
    
    def test_invalid_shape_index(self, sample_pptx_with_shapes):
        from core.powerpoint_agent_core import ShapeNotFoundError
        with pytest.raises(ShapeNotFoundError):
            format_shape(
                filepath=sample_pptx_with_shapes,
                slide_index=0,
                shape_index=99,
                fill_color="#FF0000"
            )
    
    def test_presentation_version_changes(self, sample_pptx_with_shapes):
        result = format_shape(
            filepath=sample_pptx_with_shapes,
            slide_index=0,
            shape_index=0,
            fill_color="#00FF00"
        )
        
        assert "presentation_version" in result
        assert result["presentation_version"]["before"] != result["presentation_version"]["after"]
