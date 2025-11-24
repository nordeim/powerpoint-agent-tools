import pytest
import json
import os
from pathlib import Path
from pptx import Presentation

# Assuming tools are accessible via subprocess or direct import if refactored
# For this plan, we assume direct testing of the generated files functionality

def test_add_notes(tmp_path):
    """Test adding speaker notes."""
    pptx_path = tmp_path / "test_notes.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(pptx_path)
    
    # Run tool (simulated)
    from tools.ppt_add_notes import add_notes
    result = add_notes(pptx_path, 0, "Hello World", mode="overwrite")
    
    assert result["status"] == "success"
    
    # Verify
    prs_new = Presentation(pptx_path)
    assert prs_new.slides[0].notes_slide.notes_text_frame.text == "Hello World"

def test_z_order(tmp_path):
    """Test z-order manipulation."""
    pptx_path = tmp_path / "test_z.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    
    # Add Shape 1 (Back)
    s1 = slide.shapes.add_shape(1, 0, 0, 100, 100)
    s1.name = "Shape 1"
    
    # Add Shape 2 (Front)
    s2 = slide.shapes.add_shape(1, 10, 10, 100, 100)
    s2.name = "Shape 2"
    
    prs.save(pptx_path)
    
    # Verify initial order (index)
    prs_check = Presentation(pptx_path)
    assert prs_check.slides[0].shapes[0].name == "Shape 1"
    assert prs_check.slides[0].shapes[1].name == "Shape 2"
    
    # Run tool: Send Shape 2 to back
    from tools.ppt_set_z_order import set_z_order
    # Shape 2 is at index 1
    set_z_order(pptx_path, 0, 1, "send_to_back")
    
    # Verify swap
    prs_final = Presentation(pptx_path)
    assert prs_final.slides[0].shapes[0].name == "Shape 2" # Now at back
    assert prs_final.slides[0].shapes[1].name == "Shape 1" # Now at front
