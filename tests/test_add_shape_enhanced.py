import unittest
import os
import sys
import json
import subprocess
from pptx import Presentation

class TestAddShapeEnhanced(unittest.TestCase):
    def setUp(self):
        self.test_file = "test_add_shape_enhanced.pptx"
        # Create a blank presentation
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        prs.save(self.test_file)
        self.tool_path = os.path.abspath("tools/ppt_add_shape.py")

    def tearDown(self):
        if os.path.exists(self.test_file):
            os.remove(self.test_file)

    def run_tool(self, args):
        cmd = [sys.executable, self.tool_path, "--file", self.test_file] + args
        result = subprocess.run(cmd, capture_output=True, text=True)
        return result

    def test_add_shape_with_opacity(self):
        args = [
            "--slide", "0",
            "--shape", "rectangle",
            "--position", '{"left":1.0, "top":1.0}',
            "--size", '{"width":2.0, "height":2.0}',
            "--fill-color", "#FF0000",
            "--fill-opacity", "0.5",
            "--line-color", "#0000FF",
            "--line-opacity", "0.5"
        ]
        result = self.run_tool(args)
        self.assertEqual(result.returncode, 0, f"Tool failed: {result.stderr}")
        
        output = json.loads(result.stdout)
        self.assertEqual(output["status"], "success")
        self.assertEqual(output["styling"]["fill_opacity"], 0.5)
        self.assertEqual(output["styling"]["line_opacity"], 0.5)

    def test_add_overlay(self):
        args = [
            "--slide", "0",
            "--shape", "rectangle",
            "--overlay",
            "--fill-color", "#FFFFFF"
        ]
        result = self.run_tool(args)
        self.assertEqual(result.returncode, 0, f"Tool failed: {result.stderr}")
        
        output = json.loads(result.stdout)
        self.assertEqual(output["status"], "success")
        self.assertTrue(output["is_overlay"])
        self.assertEqual(output["styling"]["fill_opacity"], 0.15) # Default overlay opacity
        
        # Check for z-order recommendation
        has_z_order_note = False
        for note in output["notes"]:
            if "ppt_set_z_order.py" in note:
                has_z_order_note = True
                break
        self.assertTrue(has_z_order_note, "Missing z-order recommendation for overlay")

if __name__ == '__main__':
    unittest.main()
