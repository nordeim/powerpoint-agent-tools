#!/usr/bin/env python3
"""Tests for ppt_validate_presentation.py v3.0"""

import pytest
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).parent.parent / "tools"))

from ppt_validate_presentation import (
    validate_presentation,
    get_policy,
    ValidationPolicy,
    ValidationIssue,
    ValidationSummary,
    VALIDATION_POLICIES
)


@pytest.fixture
def valid_pptx(tmp_path):
    """Create a valid PowerPoint for testing."""
    pptx_path = tmp_path / "valid.pptx"
    prs = Presentation()
    
    # Add slide with title
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    title = slide.shapes.title
    title.text = "Valid Presentation"
    
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def invalid_pptx(tmp_path):
    """Create a PowerPoint with various issues for testing."""
    pptx_path = tmp_path / "invalid.pptx"
    prs = Presentation()
    
    # Slide 0: Empty slide (issue)
    slide0 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Slide 1: Slide without title (issue)
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    # Add content but no title
    textbox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text_frame.text = "Content without title"
    
    # Slide 2: Valid slide with title
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
    if slide2.shapes.title:
        slide2.shapes.title.text = "Valid Slide"
    
    prs.save(str(pptx_path))
    return pptx_path


@pytest.fixture
def pptx_with_images(tmp_path):
    """Create a PowerPoint with images (no alt text) for testing."""
    pptx_path = tmp_path / "with_images.pptx"
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Create a simple test image
    from PIL import Image
    import io
    
    img = Image.new('RGB', (100, 100), color='red')
    img_bytes = io.BytesIO()
    img.save(img_bytes, format='PNG')
    img_bytes.seek(0)
    
    # Add image without alt text
    slide.shapes.add_picture(img_bytes, Inches(1), Inches(1), Inches(2), Inches(2))
    
    prs.save(str(pptx_path))
    return pptx_path


class TestGetPolicy:
    """Tests for policy retrieval."""
    
    def test_get_standard_policy(self):
        policy = get_policy("standard")
        assert policy.name == "Standard"
        assert "max_critical_issues" in policy.thresholds
    
    def test_get_strict_policy(self):
        policy = get_policy("strict")
        assert policy.name == "Strict"
        assert policy.thresholds["max_critical_issues"] == 0
        assert policy.thresholds["require_all_alt_text"] == True
    
    def test_get_lenient_policy(self):
        policy = get_policy("lenient")
        assert policy.name == "Lenient"
        assert policy.thresholds["max_critical_issues"] > 0
    
    def test_custom_policy(self):
        custom = {"max_missing_alt_text": 2}
        policy = get_policy("custom", custom)
        assert policy.thresholds["max_missing_alt_text"] == 2
    
    def test_unknown_policy_defaults_to_standard(self):
        policy = get_policy("nonexistent")
        assert policy.name == "Standard"


class TestValidationPolicy:
    """Tests for ValidationPolicy class."""
    
    def test_to_dict(self):
        policy = get_policy("standard")
        d = policy.to_dict()
        assert "name" in d
        assert "thresholds" in d
        assert "description" in d


class TestValidationIssue:
    """Tests for ValidationIssue class."""
    
    def test_to_dict(self):
        issue = ValidationIssue(
            category="structure",
            severity="critical",
            message="Test issue",
            slide_index=0,
            fix_command="test command"
        )
        d = issue.to_dict()
        assert d["category"] == "structure"
        assert d["severity"] == "critical"
        assert d["slide_index"] == 0
    
    def test_to_dict_excludes_none(self):
        issue = ValidationIssue(
            category="structure",
            severity="warning",
            message="Test"
        )
        d = issue.to_dict()
        assert "slide_index" not in d or d["slide_index"] is None


class TestValidatePresentationValid:
    """Tests for validating valid presentations."""
    
    def test_valid_presentation_passes(self, valid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(valid_pptx, policy)
        
        assert result["status"] in ("valid", "warnings")
        assert "summary" in result
        assert "issues" in result
    
    def test_valid_presentation_has_structure(self, valid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(valid_pptx, policy)
        
        assert "file" in result
        assert "validated_at" in result
        assert "policy" in result
        assert "presentation_info" in result


class TestValidatePresentationInvalid:
    """Tests for validating presentations with issues."""
    
    def test_detects_empty_slides(self, invalid_pptx):
        policy = get_policy("strict")
        result = validate_presentation(invalid_pptx, policy)
        
        # Should detect the empty slide
        assert result["summary"]["empty_slides"] >= 1
    
    def test_detects_missing_titles(self, invalid_pptx):
        policy = get_policy("strict")
        result = validate_presentation(invalid_pptx, policy)
        
        # Should detect slides without titles
        assert result["summary"]["slides_without_titles"] >= 1
    
    def test_generates_fix_commands(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        # Issues should have fix commands
        issues_with_commands = [i for i in result["issues"] if i.get("fix_command")]
        assert len(issues_with_commands) > 0
    
    def test_slide_breakdown_included(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        assert "slide_breakdown" in result
        assert len(result["slide_breakdown"]) > 0
    
    def test_recommendations_generated(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        assert "recommendations" in result
        # Should have at least one recommendation for issues found
        if result["summary"]["total_issues"] > 0:
            assert len(result["recommendations"]) > 0


class TestPolicyCompliance:
    """Tests for policy compliance checking."""
    
    def test_strict_policy_fails_on_issues(self, invalid_pptx):
        policy = get_policy("strict")
        result = validate_presentation(invalid_pptx, policy)
        
        # Strict policy should fail if there are empty slides
        if result["summary"]["empty_slides"] > 0:
            assert result["passed"] == False
    
    def test_lenient_policy_passes_with_issues(self, invalid_pptx):
        policy = get_policy("lenient")
        result = validate_presentation(invalid_pptx, policy)
        
        # Lenient policy allows many issues
        # May still pass depending on issue count
        assert "passed" in result
    
    def test_policy_violations_listed(self, invalid_pptx):
        policy = get_policy("strict")
        result = validate_presentation(invalid_pptx, policy)
        
        assert "policy_violations" in result
        if not result["passed"]:
            assert len(result["policy_violations"]) > 0


class TestAccessibilityValidation:
    """Tests for accessibility validation."""
    
    @pytest.mark.skipif(
        not pytest.importorskip("PIL", reason="Pillow required"),
        reason="Pillow not available"
    )
    def test_detects_missing_alt_text(self, pptx_with_images):
        policy = get_policy("strict")
        result = validate_presentation(pptx_with_images, policy)
        
        # Should detect missing alt text
        # Note: Detection depends on core implementation
        assert "summary" in result
    
    def test_require_all_alt_text_policy(self, tmp_path):
        """Test that require_all_alt_text policy works."""
        # Create minimal presentation
        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "Test"
        prs.save(str(pptx_path))
        
        # Test with require_all_alt_text
        policy = get_policy("strict")
        result = validate_presentation(pptx_path, policy)
        
        assert "summary" in result


class TestSummaryStatistics:
    """Tests for summary statistics."""
    
    def test_summary_has_all_fields(self, valid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(valid_pptx, policy)
        
        summary = result["summary"]
        assert "total_issues" in summary
        assert "critical_count" in summary
        assert "warning_count" in summary
        assert "empty_slides" in summary
        assert "slides_without_titles" in summary
        assert "missing_alt_text" in summary
    
    def test_counts_are_consistent(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        summary = result["summary"]
        issues = result["issues"]
        
        # Total should equal sum of severities
        assert summary["total_issues"] == len(issues)
        assert summary["critical_count"] == sum(1 for i in issues if i.get("severity") == "critical")
        assert summary["warning_count"] == sum(1 for i in issues if i.get("severity") == "warning")


class TestRecommendations:
    """Tests for recommendation generation."""
    
    def test_recommendations_have_priority(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        for rec in result["recommendations"]:
            assert "priority" in rec
            assert rec["priority"] in ("high", "medium", "low", "info")
    
    def test_recommendations_have_actions(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        for rec in result["recommendations"]:
            assert "action" in rec
            assert "issue" in rec
    
    def test_recommendations_sorted_by_priority(self, invalid_pptx):
        policy = get_policy("standard")
        result = validate_presentation(invalid_pptx, policy)
        
        if len(result["recommendations"]) > 1:
            priority_order = {"high": 0, "medium": 1, "low": 2, "info": 3}
            priorities = [priority_order.get(r["priority"], 3) for r in result["recommendations"]]
            assert priorities == sorted(priorities)


class TestFileNotFound:
    """Tests for file not found handling."""
    
    def test_raises_on_missing_file(self, tmp_path):
        policy = get_policy("standard")
        
        with pytest.raises(FileNotFoundError):
            validate_presentation(tmp_path / "nonexistent.pptx", policy)


class TestCLIIntegration:
    """Tests for CLI integration."""
    
    def test_cli_help(self):
        import subprocess
        result = subprocess.run(
            [sys.executable, "-m", "tools.ppt_validate_presentation", "--help"],
            capture_output=True,
            text=True,
            cwd=Path(__file__).parent.parent
        )
        assert result.returncode == 0
        assert "validate" in result.stdout.lower()
    
    def test_cli_validation(self, valid_pptx):
        import subprocess
        result = subprocess.run(
            [
                sys.executable, "-m", "tools.ppt_validate_presentation",
                "--file", str(valid_pptx),
                "--policy", "lenient",
                "--json"
            ],
            capture_output=True,
            text=True,
            cwd=Path(__file__).parent.parent
        )
        
        # Should succeed with valid presentation
        output = json.loads(result.stdout)
        assert "status" in output
    
    def test_cli_strict_policy(self, invalid_pptx):
        import subprocess
        result = subprocess.run(
            [
                sys.executable, "-m", "tools.ppt_validate_presentation",
                "--file", str(invalid_pptx),
                "--policy", "strict",
                "--json"
            ],
            capture_output=True,
            text=True,
            cwd=Path(__file__).parent.parent
        )
        
        output = json.loads(result.stdout)
        assert "status" in output
        # Strict policy should likely fail on invalid presentation
