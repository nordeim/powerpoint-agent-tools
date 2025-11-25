            color: Text color hex (e.g., "#FF0000")
            alignment: Text alignment ("left", "center", "right", "justify")
            
        Returns:
            Dict with shape_index and details
            
        Raises:
            SlideNotFoundError: If slide index is invalid
            InvalidPositionError: If position is invalid
        """
        slide = self._get_slide(slide_index)
        
        # Parse position and size
        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)
        
        if width is None or height is None:
            raise ValueError("Text box must have explicit width and height")
        
        # Create text box
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        
        # Configure text frame
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        # Apply formatting
        paragraph = text_frame.paragraphs[0]
        if font_name:
            paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic
        
        if color:
            paragraph.font.color.rgb = ColorHelper.from_hex(color)
        
        # Set alignment
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY
        }
        paragraph.alignment = alignment_map.get(alignment.lower(), PP_ALIGN.LEFT)
        
        # Find shape index
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "text_length": len(text),
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height}
        }
    
    def set_title(
        self,
        slide_index: int,
        title: str,
        subtitle: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Set slide title and optional subtitle.
        
        Args:
            slide_index: Target slide index
            title: Title text
            subtitle: Optional subtitle text
            
        Returns:
            Dict with title/subtitle set status
            
        Raises:
            SlideNotFoundError: If slide index is invalid
        """
        slide = self._get_slide(slide_index)
        
        title_set = False
        subtitle_set = False
        title_shape_index = None
        subtitle_shape_index = None
        
        for idx, shape in enumerate(slide.shapes):
            if shape.is_placeholder:
                ph_type = self._get_placeholder_type_int(shape.placeholder_format.type)
                
                # Check for title placeholder
                if ph_type in TITLE_PLACEHOLDER_TYPES:
                    if shape.has_text_frame:
                        shape.text_frame.text = title
                        title_set = True
                        title_shape_index = idx
                
                # Check for subtitle placeholder
                elif ph_type == SUBTITLE_PLACEHOLDER_TYPE:
                    if subtitle and shape.has_text_frame:
                        shape.text_frame.text = subtitle
                        subtitle_set = True
                        subtitle_shape_index = idx
        
        return {
            "slide_index": slide_index,
            "title_set": title_set,
            "subtitle_set": subtitle_set,
            "title_shape_index": title_shape_index,
            "subtitle_shape_index": subtitle_shape_index
        }
    
    def add_bullet_list(
        self,
        slide_index: int,
        items: List[str],
        position: Dict[str, Any],
        size: Dict[str, Any],
        bullet_style: str = "bullet",
        font_size: int = 18,
        font_name: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Add bullet list to slide.
        
        Args:
            slide_index: Target slide index
            items: List of bullet items
            position: Position dict
            size: Size dict
            bullet_style: "bullet", "numbered", or "none"
            font_size: Font size in points
            font_name: Optional font name
            
        Returns:
            Dict with shape_index and item count
        """
        slide = self._get_slide(slide_index)
        
        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)
        
        if width is None or height is None:
            raise ValueError("Bullet list must have explicit width and height")
        
        # Create text box for bullets
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for idx, item in enumerate(items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            if bullet_style == "numbered":
                p.text = f"{idx + 1}. {item}"
            else:
                p.text = item
            
            p.level = 0
            p.font.size = Pt(font_size)
            if font_name:
                p.font.name = font_name
        
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "item_count": len(items),
            "bullet_style": bullet_style
        }
    
    def format_text(
        self,
        slide_index: int,
        shape_index: int,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Format existing text shape.
        
        Args:
            slide_index: Target slide index
            shape_index: Shape index on slide
            font_name: Optional font name
            font_size: Optional font size in points
            bold: Optional bold setting
            italic: Optional italic setting
            color: Optional color hex
            
        Returns:
            Dict with formatting applied
        """
        shape = self._get_shape(slide_index, shape_index)
        
        if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
            raise ValueError(f"Shape at index {shape_index} does not have text")
        
        changes = []
        
        for paragraph in shape.text_frame.paragraphs:
            if font_name is not None:
                paragraph.font.name = font_name
                changes.append("font_name")
            if font_size is not None:
                paragraph.font.size = Pt(font_size)
                changes.append("font_size")
            if bold is not None:
                paragraph.font.bold = bold
                changes.append("bold")
            if italic is not None:
                paragraph.font.italic = italic
                changes.append("italic")
            if color is not None:
                paragraph.font.color.rgb = ColorHelper.from_hex(color)
                changes.append("color")
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "changes_applied": list(set(changes))
        }
    
    def replace_text(
        self,
        find: str,
        replace: str,
        slide_index: Optional[int] = None,
        shape_index: Optional[int] = None,
        match_case: bool = False
    ) -> Dict[str, Any]:
        """
        Find and replace text in presentation.
        
        Args:
            find: Text to find
            replace: Replacement text
            slide_index: Optional specific slide (None = all slides)
            shape_index: Optional specific shape (requires slide_index)
            match_case: Case-sensitive matching
            
        Returns:
            Dict with replacement count and locations
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        if shape_index is not None and slide_index is None:
            raise ValueError("shape_index requires slide_index to be specified")
        
        replacements = []
        total_count = 0
        
        # Determine slides to process
        if slide_index is not None:
            slides_to_process = [(slide_index, self._get_slide(slide_index))]
        else:
            slides_to_process = list(enumerate(self.prs.slides))
        
        for s_idx, slide in slides_to_process:
            # Determine shapes to process
            if shape_index is not None:
                shapes_to_process = [(shape_index, self._get_shape(s_idx, shape_index))]
            else:
                shapes_to_process = list(enumerate(slide.shapes))
            
            for sh_idx, shape in shapes_to_process:
                if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
                    continue
                
                count = self._replace_text_in_shape(shape, find, replace, match_case)
                if count > 0:
                    total_count += count
                    replacements.append({
                        "slide": s_idx,
                        "shape": sh_idx,
                        "count": count
                    })
        
        return {
            "find": find,
            "replace": replace,
            "match_case": match_case,
            "total_replacements": total_count,
            "locations": replacements
        }
    
    def _replace_text_in_shape(
        self,
        shape,
        find: str,
        replace: str,
        match_case: bool
    ) -> int:
        """Replace text within a single shape, preserving formatting where possible."""
        count = 0
        
        try:
            text_frame = shape.text_frame
        except (AttributeError, TypeError):
            return 0
        
        # Strategy 1: Replace in runs (preserves formatting)
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if match_case:
                    if find in run.text:
                        occurrences = run.text.count(find)
                        run.text = run.text.replace(find, replace)
                        count += occurrences
                else:
                    if find.lower() in run.text.lower():
                        pattern = re.compile(re.escape(find), re.IGNORECASE)
                        matches = pattern.findall(run.text)
                        run.text = pattern.sub(replace, run.text)
                        count += len(matches)
        
        if count > 0:
            return count
        
        # Strategy 2: Full text replacement (if text spans runs)
        try:
            full_text = shape.text
            if not full_text:
                return 0
            
            if match_case:
                if find in full_text:
                    occurrences = full_text.count(find)
                    shape.text = full_text.replace(find, replace)
                    return occurrences
            else:
                if find.lower() in full_text.lower():
                    pattern = re.compile(re.escape(find), re.IGNORECASE)
                    matches = pattern.findall(full_text)
                    shape.text = pattern.sub(replace, full_text)
                    return len(matches)
        except (AttributeError, TypeError):
            pass
        
        return 0
    
    def add_notes(
        self,
        slide_index: int,
        text: str,
        mode: str = "append"
    ) -> Dict[str, Any]:
        """
        Add speaker notes to a slide.
        
        Args:
            slide_index: Target slide index
            text: Notes text to add
            mode: "append", "prepend", or "overwrite"
            
        Returns:
            Dict with notes details
            
        Raises:
            SlideNotFoundError: If slide index is invalid
            ValueError: If mode is invalid
        """
        if mode not in ("append", "prepend", "overwrite"):
            raise ValueError(f"Invalid mode: {mode}. Must be 'append', 'prepend', or 'overwrite'")
        
        slide = self._get_slide(slide_index)
        
        # Access or create notes slide
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        
        original_text = text_frame.text or ""
        original_length = len(original_text)
        
        if mode == "overwrite":
            text_frame.text = text
            final_text = text
        elif mode == "append":
            if original_text.strip():
                final_text = original_text + "\n" + text
            else:
                final_text = text
            text_frame.text = final_text
        elif mode == "prepend":
            if original_text.strip():
                final_text = text + "\n" + original_text
            else:
                final_text = text
            text_frame.text = final_text
        
        return {
            "slide_index": slide_index,
            "mode": mode,
            "original_length": original_length,
            "new_length": len(final_text),
            "text_preview": final_text[:100] + "..." if len(final_text) > 100 else final_text
        }
    
    def set_footer(
        self,
        text: Optional[str] = None,
        show_slide_number: bool = False,
        show_date: bool = False,
        slide_index: Optional[int] = None
    ) -> Dict[str, Any]:
        """
        Set footer properties for slide(s).
        
        Note: Footer configuration in python-pptx is limited.
        This method sets footer placeholders where available.
        
        Args:
            text: Footer text
            show_slide_number: Show slide numbers
            show_date: Show date
            slide_index: Specific slide (None = all slides)
            
        Returns:
            Dict with footer configuration results
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        results = []
        
        # Determine slides to process
        if slide_index is not None:
            slides = [(slide_index, self._get_slide(slide_index))]
        else:
            slides = list(enumerate(self.prs.slides))
        
        for s_idx, slide in slides:
            slide_result = {
                "slide_index": s_idx,
                "footer_set": False,
                "slide_number_set": False,
                "date_set": False
            }
            
            for shape in slide.shapes:
                if not shape.is_placeholder:
                    continue
                
                ph_type = self._get_placeholder_type_int(shape.placeholder_format.type)
                
                # Footer placeholder (type 7)
                if ph_type == 7 and text is not None:
                    if shape.has_text_frame:
                        shape.text_frame.text = text
                        slide_result["footer_set"] = True
                
                # Slide number placeholder (type 6)
                if ph_type == 6 and show_slide_number:
                    slide_result["slide_number_set"] = True
                
                # Date placeholder (type 5)
                if ph_type == 5 and show_date:
                    slide_result["date_set"] = True
            
            results.append(slide_result)
        
        return {
            "text": text,
            "show_slide_number": show_slide_number,
            "show_date": show_date,
            "slides_processed": len(results),
            "results": results
        }
    
    # ========================================================================
    # SHAPE OPERATIONS
    # ========================================================================
    
    def add_shape(
        self,
        slide_index: int,
        shape_type: str,
        position: Dict[str, Any],
        size: Dict[str, Any],
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
        line_width: float = 1.0,
        text: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Add shape to slide.
        
        Args:
            slide_index: Target slide index
            shape_type: Shape type name (rectangle, ellipse, arrow_right, etc.)
            position: Position dict
            size: Size dict
            fill_color: Fill color hex
            line_color: Line color hex
            line_width: Line width in points
            text: Optional text to add inside shape
            
        Returns:
            Dict with shape_index and details
        """
        slide = self._get_slide(slide_index)
        
        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)
        
        if width is None or height is None:
            raise ValueError("Shape must have explicit width and height")
        
        # Map shape type string to MSO constant
        shape_type_map = {
            "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
            "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
            "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
            "arrow_right": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            "arrow_left": MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
            "arrow_up": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
            "arrow_down": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
            "pentagon": MSO_AUTO_SHAPE_TYPE.PENTAGON,
            "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
            "star": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            "heart": MSO_AUTO_SHAPE_TYPE.HEART,
            "lightning": MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT,
            "sun": MSO_AUTO_SHAPE_TYPE.SUN,
            "moon": MSO_AUTO_SHAPE_TYPE.MOON,
            "cloud": MSO_AUTO_SHAPE_TYPE.CLOUD,
        }
        
        mso_shape = shape_type_map.get(
            shape_type.lower(),
            MSO_AUTO_SHAPE_TYPE.RECTANGLE
        )
        
        # Add shape
        shape = slide.shapes.add_shape(
            mso_shape,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        
        # Apply fill color
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = ColorHelper.from_hex(fill_color)
        
        # Apply line color
        if line_color:
            shape.line.color.rgb = ColorHelper.from_hex(line_color)
            shape.line.width = Pt(line_width)
        
        # Add text if provided
        if text and shape.has_text_frame:
            shape.text_frame.text = text
        
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "shape_type": shape_type,
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height}
        }
    
    def format_shape(
        self,
        slide_index: int,
        shape_index: int,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
        line_width: Optional[float] = None,
        transparency: Optional[float] = None
    ) -> Dict[str, Any]:
        """
        Format existing shape.
        
        Args:
            slide_index: Target slide index
            shape_index: Shape index on slide
            fill_color: Fill color hex
            line_color: Line color hex
            line_width: Line width in points
            transparency: Fill transparency (0.0 = opaque, 1.0 = invisible)
            
        Returns:
            Dict with formatting changes applied
        """
        shape = self._get_shape(slide_index, shape_index)
        
        changes = []
        
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = ColorHelper.from_hex(fill_color)
            changes.append("fill_color")
        
        if line_color is not None:
            shape.line.color.rgb = ColorHelper.from_hex(line_color)
            changes.append("line_color")
        
        if line_width is not None:
            shape.line.width = Pt(line_width)
            changes.append("line_width")
        
        if transparency is not None:
            try:
                # Transparency is set on fill
                shape.fill.solid()
                # Note: python-pptx doesn't directly support transparency
                # This is a best-effort implementation
                changes.append("transparency_attempted")
            except Exception:
                pass
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "changes_applied": changes
        }
    
    def remove_shape(self, slide_index: int, shape_index: int) -> Dict[str, Any]:
        """
        Remove shape from slide.
        
        Args:
            slide_index: Target slide index
            shape_index: Shape index to remove
            
        Returns:
            Dict with removal details
            
        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
        """
        slide = self._get_slide(slide_index)
        shape = self._get_shape(slide_index, shape_index)
        
        # Get shape info before removal
        shape_name = shape.name
        shape_type = str(shape.shape_type)
        
        # Remove shape from slide
        sp = shape.element
        sp.getparent().remove(sp)
        
        return {
            "slide_index": slide_index,
            "removed_shape_index": shape_index,
            "removed_shape_name": shape_name,
            "removed_shape_type": shape_type,
            "new_shape_count": len(slide.shapes)
        }
    
    def set_z_order(
        self,
        slide_index: int,
        shape_index: int,
        action: str
    ) -> Dict[str, Any]:
        """
        Change the z-order (stacking order) of a shape.
        
        Args:
            slide_index: Target slide index
            shape_index: Shape index to modify
            action: One of "bring_to_front", "send_to_back", 
                   "bring_forward", "send_backward"
            
        Returns:
            Dict with z-order change details including old and new positions
            
        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
            ValueError: If action is invalid
        """
        valid_actions = {"bring_to_front", "send_to_back", "bring_forward", "send_backward"}
        if action not in valid_actions:
            raise ValueError(f"Invalid action: {action}. Must be one of {valid_actions}")
        
        slide = self._get_slide(slide_index)
        shape = self._get_shape(slide_index, shape_index)
        
        # Access the shape tree XML element
        sp_tree = slide.shapes._spTree
        element = shape.element
        
        # Find current position in XML tree
        current_index = -1
        shape_elements = [child for child in sp_tree if child.tag.endswith('}sp') or 
                         child.tag.endswith('}pic') or child.tag.endswith('}graphicFrame')]
        
        for i, child in enumerate(sp_tree):
            if child == element:
                current_index = i
                break
        
        if current_index == -1:
            raise PowerPointAgentError(
                "Could not locate shape in XML tree",
                details={"slide_index": slide_index, "shape_index": shape_index}
            )
        
        new_index = current_index
        max_index = len(sp_tree) - 1
        
        # Execute the z-order action
        if action == "bring_to_front":
            sp_tree.remove(element)
            sp_tree.append(element)
            new_index = len(sp_tree) - 1
            
        elif action == "send_to_back":
            sp_tree.remove(element)
            # Insert after nvGrpSpPr and grpSpPr (indices 0 and 1 typically)
            sp_tree.insert(2, element)
            new_index = 2
            
        elif action == "bring_forward":
            if current_index < max_index:
                sp_tree.remove(element)
                sp_tree.insert(current_index + 1, element)
                new_index = current_index + 1
                
        elif action == "send_backward":
            if current_index > 2:  # Don't go before required elements
                sp_tree.remove(element)
                sp_tree.insert(current_index - 1, element)
                new_index = current_index - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "action": action,
            "z_order_change": {
                "from": current_index,
                "to": new_index
            },
            "warning": "Shape indices may have changed after z-order operation. Re-query slide info."
        }
    
    def add_table(
        self,
        slide_index: int,
        rows: int,
        cols: int,
        position: Dict[str, Any],
        size: Dict[str, Any],
        data: Optional[List[List[Any]]] = None,
        header_row: bool = True
    ) -> Dict[str, Any]:
        """
        Add table to slide.
        
        Args:
            slide_index: Target slide index
            rows: Number of rows
            cols: Number of columns
            position: Position dict
            size: Size dict
            data: Optional 2D list of cell values
            header_row: Whether first row is header (styling hint)
            
        Returns:
            Dict with shape_index and table details
        """
        slide = self._get_slide(slide_index)
        
        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)
        
        if width is None or height is None:
            raise ValueError("Table must have explicit width and height")
        
        # Create table
        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        
        table = table_shape.table
        
        # Populate with data if provided
        cells_filled = 0
        if data:
            for row_idx, row_data in enumerate(data):
                if row_idx >= rows:
                    break
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx >= cols:
                        break
                    table.cell(row_idx, col_idx).text = str(cell_value)
                    cells_filled += 1
        
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "rows": rows,
            "cols": cols,
            "cells_filled": cells_filled,
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height}
        }
    
    def add_connector(
        self,
        slide_index: int,
        from_shape_index: int,
        to_shape_index: int,
        connector_type: str = "straight"
    ) -> Dict[str, Any]:
        """
        Add connector line between two shapes.
        
        Args:
            slide_index: Target slide index
            from_shape_index: Starting shape index
            to_shape_index: Ending shape index
            connector_type: "straight", "elbow", or "curved"
            
        Returns:
            Dict with connector details
        """
        slide = self._get_slide(slide_index)
        
        shape1 = self._get_shape(slide_index, from_shape_index)
        shape2 = self._get_shape(slide_index, to_shape_index)
        
        # Calculate center points
        x1 = shape1.left + shape1.width // 2
        y1 = shape1.top + shape1.height // 2
        x2 = shape2.left + shape2.width // 2
        y2 = shape2.top + shape2.height // 2
        
        # Map connector type
        connector_map = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curved": MSO_CONNECTOR.CURVE
        }
        mso_connector = connector_map.get(connector_type.lower(), MSO_CONNECTOR.STRAIGHT)
        
        # Add connector
        connector = slide.shapes.add_connector(
            mso_connector,
            x1, y1, x2, y2
        )
        
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "from_shape": from_shape_index,
            "to_shape": to_shape_index,
            "connector_type": connector_type
        }
    
    # ========================================================================
    # IMAGE OPERATIONS
    # ========================================================================
    
    def insert_image(
        self,
        slide_index: int,
        image_path: Union[str, Path],
        position: Dict[str, Any],
        size: Optional[Dict[str, Any]] = None,
        alt_text: Optional[str] = None,
        compress: bool = False
    ) -> Dict[str, Any]:
        """
        Insert image on slide.
        
        Args:
            slide_index: Target slide index
            image_path: Path to image file
            position: Position dict
            size: Optional size dict (can use "auto" for aspect ratio)
            alt_text: Alternative text for accessibility
            compress: Compress image before inserting
            
        Returns:
            Dict with shape_index and image details
        """
        slide = self._get_slide(slide_index)
        image_path = PathValidator.validate_image_path(image_path)
        
        left, top = Position.from_dict(position)
        
        # Get aspect ratio if Pillow available
        aspect_ratio = None
        if HAS_PILLOW:
            try:
                with PILImage.open(image_path) as img:
                    aspect_ratio = img.width / img.height
            except Exception:
                pass
        
        # Parse size
        if size:
            width, height = Size.from_dict(size, aspect_ratio=aspect_ratio)
        else:
            # Default to half slide width, maintain aspect ratio
            width = SLIDE_WIDTH_INCHES * 0.5
            if aspect_ratio:
                height = width / aspect_ratio
            else:
                height = SLIDE_HEIGHT_INCHES * 0.3
        
        # Compress if requested
        if compress and HAS_PILLOW:
            image_stream = AssetValidator.compress_image(image_path)
            picture = slide.shapes.add_picture(
                image_stream,
                Inches(left), Inches(top),
                width=Inches(width) if width else None,
                height=Inches(height) if height else None
            )
        else:
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(left), Inches(top),
                width=Inches(width) if width else None,
                height=Inches(height) if height else None
            )
        
        # Set alt text
        if alt_text:
            picture.name = alt_text
            try:
                # Set description attribute for proper alt text
                picture._element.set('descr', alt_text)
            except Exception:
                pass
        
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "image_path": str(image_path),
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height},
            "alt_text_set": alt_text is not None,
            "compressed": compress
        }
    
    def replace_image(
        self,
        slide_index: int,
        old_image_name: str,
        new_image_path: Union[str, Path],
        compress: bool = False
    ) -> Dict[str, Any]:
        """
        Replace existing image by name.
        
        Args:
            slide_index: Target slide index
            old_image_name: Name or partial name of image to replace
            new_image_path: Path to new image file
            compress: Compress new image
            
        Returns:
            Dict with replacement details
        """
        slide = self._get_slide(slide_index)
        new_image_path = PathValidator.validate_image_path(new_image_path)
        
        replaced = False
        old_shape_index = None
        new_shape_index = None
        
        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.name == old_image_name or old_image_name in (shape.name or ""):
                    # Store position and size
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    old_shape_index = idx
                    
                    # Remove old image
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
                    # Add new image
                    if compress and HAS_PILLOW:
                        image_stream = AssetValidator.compress_image(new_image_path)
                        new_picture = slide.shapes.add_picture(
                            image_stream, left, top,
                            width=width, height=height
                        )
                    else:
                        new_picture = slide.shapes.add_picture(
                            str(new_image_path), left, top,
                            width=width, height=height
                        )
                    
                    new_shape_index = len(slide.shapes) - 1
                    replaced = True
                    break
        
        return {
            "slide_index": slide_index,
            "replaced": replaced,
            "old_image_name": old_image_name,
            "old_shape_index": old_shape_index,
            "new_image_path": str(new_image_path),
            "new_shape_index": new_shape_index
        }
    
    def set_image_properties(
        self,
        slide_index: int,
        shape_index: int,
        alt_text: Optional[str] = None,
        name: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Set image properties.
        
        Args:
            slide_index: Target slide index
            shape_index: Image shape index
            alt_text: Alternative text for accessibility
            name: Shape name
            
        Returns:
            Dict with properties set
        """
        shape = self._get_shape(slide_index, shape_index)
        
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape at index {shape_index} is not an image")
        
        changes = []
        
        if alt_text is not None:
            try:
                shape._element.set('descr', alt_text)
                changes.append("alt_text")
            except Exception:
                # Fallback to name
                shape.name = alt_text
                changes.append("alt_text_via_name")
        
        if name is not None:
            shape.name = name
            changes.append("name")
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "changes_applied": changes
        }
    
    def crop_image(
        self,
        slide_index: int,
        shape_index: int,
        left: float = 0.0,
        top: float = 0.0,
        right: float = 0.0,
        bottom: float = 0.0
    ) -> Dict[str, Any]:
        """
        Crop image by specifying crop amounts from each edge.
        
        Args:
            slide_index: Target slide index
            shape_index: Image shape index
            left: Crop from left (0.0 to 1.0, proportion of width)
            top: Crop from top (0.0 to 1.0, proportion of height)
            right: Crop from right (0.0 to 1.0, proportion of width)
            bottom: Crop from bottom (0.0 to 1.0, proportion of height)
            
        Returns:
            Dict with crop details
        """
        shape = self._get_shape(slide_index, shape_index)
        
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape at index {shape_index} is not an image")
        
        # Validate crop values
        for name, value in [("left", left), ("top", top), ("right", right), ("bottom", bottom)]:
            if not 0.0 <= value < 1.0:
                raise ValueError(f"Crop {name} must be between 0.0 and 1.0, got {value}")
        
        if left + right >= 1.0:
            raise ValueError("Left + right crop cannot equal or exceed 1.0")
        if top + bottom >= 1.0:
            raise ValueError("Top + bottom crop cannot equal or exceed 1.0")
        
        # Apply crop using picture's crop properties
        try:
            # Access the picture element
            pic = shape._element
            
            # Find or create blipFill element
            blip_fill = pic.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill')
            if blip_fill is None:
                blip_fill = pic.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill')
            
            if blip_fill is not None:
                # Find or create srcRect element
                ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                src_rect = blip_fill.find(f'{ns}srcRect')
                
                if src_rect is None:
                    from lxml import etree
                    src_rect = etree.SubElement(blip_fill, f'{ns}srcRect')
                
                # Set crop values (in percentage * 1000)
                src_rect.set('l', str(int(left * 100000)))
                src_rect.set('t', str(int(top * 100000)))
                src_rect.set('r', str(int(right * 100000)))
                src_rect.set('b', str(int(bottom * 100000)))
                
                return {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "crop_applied": True,
                    "crop_values": {
                        "left": left,
                        "top": top,
                        "right": right,
                        "bottom": bottom
                    }
                }
        except Exception as e:
            logger.warning(f"Could not apply crop via XML: {e}")
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "crop_applied": False,
            "error": "Crop not supported for this image type"
        }
    
    def resize_image(
        self,
        slide_index: int,
        shape_index: int,
        width: Optional[float] = None,
        height: Optional[float] = None,
        maintain_aspect: bool = True
    ) -> Dict[str, Any]:
        """
        Resize image shape.
        
        Args:
            slide_index: Target slide index
            shape_index: Image shape index
            width: New width in inches (None = keep current)
            height: New height in inches (None = keep current)
            maintain_aspect: Maintain aspect ratio
            
        Returns:
            Dict with new dimensions
        """
        shape = self._get_shape(slide_index, shape_index)
        
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape at index {shape_index} is not an image")
        
        original_width = shape.width / EMU_PER_INCH
        original_height = shape.height / EMU_PER_INCH
        aspect = original_width / original_height if original_height > 0 else 1.0
        
        new_width = width
        new_height = height
        
        if maintain_aspect:
            if width is not None and height is None:
                new_height = width / aspect
            elif height is not None and width is None:
                new_width = height * aspect
        
        if new_width is not None:
            shape.width = Inches(new_width)
        if new_height is not None:
            shape.height = Inches(new_height)
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "original_size": {"width": original_width, "height": original_height},
            "new_size": {
                "width": new_width or original_width,
                "height": new_height or original_height
            }
        }
    
    # ========================================================================
    # CHART OPERATIONS
    # ========================================================================
    
    def add_chart(
        self,
        slide_index: int,
        chart_type: str,
        data: Dict[str, Any],
        position: Dict[str, Any],
        size: Dict[str, Any],
        title: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Add chart to slide.
        
        Args:
            slide_index: Target slide index
            chart_type: Chart type (column, bar, line, pie, etc.)
            data: Chart data dict with "categories" and "series"
            position: Position dict
            size: Size dict
            title: Optional chart title
            
        Returns:
            Dict with shape_index and chart details
            
        Example data:
            {
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "series": [
                    {"name": "Revenue", "values": [100, 120, 140, 160]},
                    {"name": "Costs", "values": [80, 90, 100, 110]}
                ]
            }
        """
        slide = self._get_slide(slide_index)
        
        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)
        
        if width is None or height is None:
            raise ValueError("Chart must have explicit width and height")
        
        # Map chart type string to XL constant
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
            "area": XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }
        
        xl_chart_type = chart_type_map.get(
            chart_type.lower(),
            XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        
        # Build chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])
        
        for series in data.get("series", []):
            chart_data.add_series(series["name"], series["values"])
        
        # Add chart
        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        )
        
        # Set title if provided
        if title:
            chart_shape.chart.has_title = True
            chart_shape.chart.chart_title.text_frame.text = title
        
        shape_index = len(slide.shapes) - 1
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "chart_type": chart_type,
            "categories_count": len(data.get("categories", [])),
            "series_count": len(data.get("series", [])),
            "title": title,
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height}
        }
    
    def update_chart_data(
        self,
        slide_index: int,
        chart_index: int,
        data: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Update existing chart data.
        
        Args:
            slide_index: Target slide index
            chart_index: Chart index on slide (not shape index)
            data: New chart data dict
            
        Returns:
            Dict with update details
        """
        chart_shape = self._get_chart_shape(slide_index, chart_index)
        
        # Build new chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])
        
        for series in data.get("series", []):
            chart_data.add_series(series["name"], series["values"])
        
        # Try to replace data (preserves formatting)
        try:
            chart_shape.chart.replace_data(chart_data)
            method = "replace_data"
        except AttributeError:
            # Fallback: recreate chart (loses some formatting)
            logger.warning(
                "chart.replace_data() not available. "
                "Recreating chart (some formatting may be lost)."
            )
            
            slide = self._get_slide(slide_index)
            
            # Store chart properties
            left = chart_shape.left
            top = chart_shape.top
            width = chart_shape.width
            height = chart_shape.height
            chart_type = chart_shape.chart.chart_type
            has_title = chart_shape.chart.has_title
            title_text = None
            if has_title:
                try:
                    title_text = chart_shape.chart.chart_title.text_frame.text
                except Exception:
                    pass
            
            # Remove old chart
            sp = chart_shape.element
            sp.getparent().remove(sp)
            
            # Create new chart
            new_chart_shape = slide.shapes.add_chart(
                chart_type, left, top, width, height, chart_data
            )
            
            # Restore title
            if title_text:
                new_chart_shape.chart.has_title = True
                new_chart_shape.chart.chart_title.text_frame.text = title_text
            
            method = "recreate"
        
        return {
            "slide_index": slide_index,
            "chart_index": chart_index,
            "categories_count": len(data.get("categories", [])),
            "series_count": len(data.get("series", [])),
            "update_method": method
        }
    
    def format_chart(
        self,
        slide_index: int,
        chart_index: int,
        title: Optional[str] = None,
        legend_position: Optional[str] = None,
        has_legend: Optional[bool] = None
    ) -> Dict[str, Any]:
        """
        Format existing chart.
        
        Args:
            slide_index: Target slide index
            chart_index: Chart index on slide
            title: Chart title
            legend_position: Legend position ("bottom", "left", "right", "top")
            has_legend: Show/hide legend
            
        Returns:
            Dict with formatting changes
        """
        chart_shape = self._get_chart_shape(slide_index, chart_index)
        chart = chart_shape.chart
        
        changes = []
        
        if title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            changes.append("title")
        
        if has_legend is not None:
            chart.has_legend = has_legend
            changes.append("has_legend")
        
        if legend_position is not None and chart.has_legend:
            from pptx.enum.chart import XL_LEGEND_POSITION
            position_map = {
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "left": XL_LEGEND_POSITION.LEFT,
                "right": XL_LEGEND_POSITION.RIGHT,
                "top": XL_LEGEND_POSITION.TOP,
                "corner": XL_LEGEND_POSITION.CORNER,
            }
            if legend_position.lower() in position_map:
                chart.legend.position = position_map[legend_position.lower()]
                changes.append("legend_position")
        
        return {
            "slide_index": slide_index,
            "chart_index": chart_index,
            "changes_applied": changes
        }
    
    # ========================================================================
    # LAYOUT & THEME OPERATIONS
    # ========================================================================
    
    def set_slide_layout(self, slide_index: int, layout_name: str) -> Dict[str, Any]:
        """
        Change slide layout.
        
        Note: This changes the layout but may not reposition existing content.
        
        Args:
            slide_index: Target slide index
            layout_name: Name of new layout
            
        Returns:
            Dict with layout change details
        """
        slide = self._get_slide(slide_index)
        layout = self._get_layout(layout_name)
        
        old_layout = slide.slide_layout.name
        slide.slide_layout = layout
        
        return {
            "slide_index": slide_index,
            "old_layout": old_layout,
            "new_layout": layout_name
        }
    
    def set_background(
        self,
        slide_index: Optional[int] = None,
        color: Optional[str] = None,
        image_path: Optional[Union[str, Path]] = None
    ) -> Dict[str, Any]:
        """
        Set slide background color or image.
        
        Args:
            slide_index: Target slide (None = all slides)
            color: Background color hex
            image_path: Background image path
            
        Returns:
            Dict with background change details
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        if color is None and image_path is None:
            raise ValueError("Must specify either color or image_path")
        
        results = []
        
        # Determine slides to process
        if slide_index is not None:
            slides = [(slide_index, self._get_slide(slide_index))]
        else:
            slides = list(enumerate(self.prs.slides))
        
        for s_idx, slide in slides:
            result = {"slide_index": s_idx, "success": False}
            
            try:
                background = slide.background
                fill = background.fill
                
                if color:
                    fill.solid()
                    fill.fore_color.rgb = ColorHelper.from_hex(color)
                    result["success"] = True
                    result["type"] = "color"
                    result["color"] = color
                
                elif image_path:
                    # Note: python-pptx has limited background image support
                    # This is a best-effort implementation
                    image_path = PathValidator.validate_image_path(image_path)
                    result["type"] = "image"
                    result["image_path"] = str(image_path)
                    result["note"] = "Background image support is limited in python-pptx"
                    
            except Exception as e:
                result["error"] = str(e)
            
            results.append(result)
        
        return {
            "slides_processed": len(results),
            "results": results
        }
    
    def get_available_layouts(self) -> List[str]:
        """
        Get list of available layout names.
        
        Returns:
            List of layout name strings
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        self._ensure_layout_cache()
        return list(self._layout_cache.keys())
    
    # ========================================================================
    # VALIDATION OPERATIONS
    # ========================================================================
    
    def validate_presentation(self) -> Dict[str, Any]:
        """
        Comprehensive presentation validation.
        
        Returns:
            Validation report dict
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        issues = {
            "empty_slides": [],
            "slides_without_titles": [],
            "fonts_used": set(),
            "large_shapes": []
        }
        
        for idx, slide in enumerate(self.prs.slides):
            # Check for empty slides
            if len(slide.shapes) == 0:
                issues["empty_slides"].append(idx)
            
            # Check for title
            has_title = False
            for shape in slide.shapes:
                if shape.is_placeholder:
                    ph_type = self._get_placeholder_type_int(shape.placeholder_format.type)
                    if ph_type in TITLE_PLACEHOLDER_TYPES:
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            has_title = True
                            break
                
                # Collect fonts
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.font.name:
                            issues["fonts_used"].add(para.font.name)
            
            if not has_title:
                issues["slides_without_titles"].append(idx)
        
        issues["fonts_used"] = list(issues["fonts_used"])
        
        total_issues = (
            len(issues["empty_slides"]) +
            len(issues["slides_without_titles"])
        )
        
        return {
            "status": "issues_found" if total_issues > 0 else "valid",
            "total_issues": total_issues,
            "slide_count": len(self.prs.slides),
            "issues": issues
        }
    
    def check_accessibility(self) -> Dict[str, Any]:
        """
        Run accessibility checker.
        
        Returns:
            Accessibility report dict
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        return AccessibilityChecker.check_presentation(self.prs)
    
    def validate_assets(self) -> Dict[str, Any]:
        """
        Run asset validator.
        
        Returns:
            Asset validation report dict
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        return AssetValidator.validate_presentation_assets(self.prs, self.filepath)
    
    # ========================================================================
    # EXPORT OPERATIONS
    # ========================================================================
    
    def export_to_pdf(self, output_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Export presentation to PDF.
        
        Requires LibreOffice or Microsoft Office installed.
        
        Args:
            output_path: Output PDF path
            
        Returns:
            Dict with export details
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        output_path = Path(output_path)
        if output_path.suffix.lower() != '.pdf':
            output_path = output_path.with_suffix('.pdf')
        
        # Ensure parent directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Save to temp file first
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            temp_pptx = Path(tmp.name)
        
        try:
            self.prs.save(str(temp_pptx))
            
            # Try LibreOffice conversion
            result = subprocess.run(
                [
                    'soffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', str(output_path.parent), str(temp_pptx)
                ],
                capture_output=True,
                timeout=120
            )
            
            if result.returncode != 0:
                raise PowerPointAgentError(
                    "PDF export failed. LibreOffice is required for PDF export.",
                    details={
                        "stderr": result.stderr.decode() if result.stderr else None,
                        "install_instructions": {
                            "linux": "sudo apt install libreoffice-impress",
                            "macos": "brew install --cask libreoffice",
                            "windows": "Download from libreoffice.org"
                        }
                    }
                )
            
            # Rename output file to desired name
            generated_pdf = output_path.parent / f"{temp_pptx.stem}.pdf"
            if generated_pdf.exists() and generated_pdf != output_path:
                shutil.move(str(generated_pdf), str(output_path))
            
            return {
                "success": True,
                "output_path": str(output_path),
                "file_size_bytes": output_path.stat().st_size if output_path.exists() else 0
            }
            
        finally:
            temp_pptx.unlink(missing_ok=True)
    
    def extract_notes(self) -> Dict[int, str]:
        """
        Extract speaker notes from all slides.
        
        Returns:
            Dict mapping slide index to notes text
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        notes = {}
        
        for idx, slide in enumerate(self.prs.slides):
            if slide.has_notes_slide:
                try:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    if text_frame.text and text_frame.text.strip():
                        notes[idx] = text_frame.text
                except Exception:
                    pass
        
        return notes
    
    # ========================================================================
    # INFORMATION & VERSIONING
    # ========================================================================
    
    def get_presentation_info(self) -> Dict[str, Any]:
        """
        Get presentation metadata and information.
        
        Returns:
            Dict with presentation information
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        info = {
            "slide_count": len(self.prs.slides),
            "layouts": self.get_available_layouts(),
            "slide_width_inches": self.prs.slide_width / EMU_PER_INCH,
            "slide_height_inches": self.prs.slide_height / EMU_PER_INCH,
            "presentation_version": self.get_presentation_version()
        }
        
        # Calculate aspect ratio
        width = info["slide_width_inches"]
        height = info["slide_height_inches"]
        if height > 0:
            ratio = width / height
            if abs(ratio - 16/9) < 0.1:
                info["aspect_ratio"] = "16:9"
            elif abs(ratio - 4/3) < 0.1:
                info["aspect_ratio"] = "4:3"
            else:
                info["aspect_ratio"] = f"{width:.2f}:{height:.2f}"
        
        # File info
        if self.filepath and self.filepath.exists():
            stat = self.filepath.stat()
            info["file"] = str(self.filepath)
            info["file_size_bytes"] = stat.st_size
            info["file_size_mb"] = round(stat.st_size / (1024 * 1024), 2)
            info["modified"] = datetime.fromtimestamp(stat.st_mtime).isoformat()
        
        return info
    
    def get_slide_info(self, slide_index: int) -> Dict[str, Any]:
        """
        Get detailed information about a specific slide.
        
        Args:
            slide_index: Slide index to inspect
            
        Returns:
            Dict with comprehensive slide information
        """
        slide = self._get_slide(slide_index)
        
        shapes_info = []
        for idx, shape in enumerate(slide.shapes):
            # Determine shape type string
            shape_type_str = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
            
            if shape.is_placeholder:
                ph_type = self._get_placeholder_type_int(shape.placeholder_format.type)
                ph_name = get_placeholder_type_name(ph_type)
                shape_type_str = f"PLACEHOLDER ({ph_name})"
            
            shape_info = {
                "index": idx,
                "type": shape_type_str,
                "name": shape.name,
                "has_text": hasattr(shape, 'text_frame') and shape.has_text_frame,
                "position": {
                    "left_inches": round(shape.left / EMU_PER_INCH, 3),
                    "top_inches": round(shape.top / EMU_PER_INCH, 3),
                    "left_percent": f"{(shape.left / self.prs.slide_width * 100):.1f}%",
                    "top_percent": f"{(shape.top / self.prs.slide_height * 100):.1f}%"
                },
                "size": {
                    "width_inches": round(shape.width / EMU_PER_INCH, 3),
                    "height_inches": round(shape.height / EMU_PER_INCH, 3),
                    "width_percent": f"{(shape.width / self.prs.slide_width * 100):.1f}%",
                    "height_percent": f"{(shape.height / self.prs.slide_height * 100):.1f}%"
                }
            }
            
            # Add text content if present
            if shape.has_text_frame:
                try:
                    full_text = shape.text_frame.text
                    shape_info["text"] = full_text
                    shape_info["text_length"] = len(full_text)
                except Exception:
                    pass
            
            # Add image info if picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    shape_info["image_size_bytes"] = len(shape.image.blob)
                    shape_info["image_content_type"] = shape.image.content_type
                except Exception:
                    pass
            
            # Add chart info if chart
            if hasattr(shape, 'has_chart') and shape.has_chart:
                try:
                    shape_info["chart_type"] = str(shape.chart.chart_type)
                except Exception:
                    pass
            
            shapes_info.append(shape_info)
        
        # Check for notes
        has_notes = False
        notes_preview = None
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text and notes_text.strip():
                    has_notes = True
                    notes_preview = notes_text[:100] + "..." if len(notes_text) > 100 else notes_text
            except Exception:
                pass
        
        return {
            "slide_index": slide_index,
            "layout": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "shapes": shapes_info,
            "has_notes": has_notes,
            "notes_preview": notes_preview
        }
    
    def get_presentation_version(self) -> str:
        """
        Compute a deterministic version hash for the presentation.
        
        The version is based on:
        - Slide count
        - Layout names
        - Shape counts per slide
        - Text content hashes
        
        Returns:
            SHA-256 hash prefix (16 characters)
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        # Build version components
        components = []
        
        # Slide count
        components.append(f"slides:{len(self.prs.slides)}")
        
        # Per-slide information
        for idx, slide in enumerate(self.prs.slides):
            slide_components = [
                f"slide:{idx}",
                f"layout:{slide.slide_layout.name}",
                f"shapes:{len(slide.shapes)}"
            ]
            
            # Add text content hash
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    try:
                        text_content.append(shape.text_frame.text)
                    except Exception:
                        pass
            
            if text_content:
                text_hash = hashlib.md5("".join(text_content).encode()).hexdigest()[:8]
                slide_components.append(f"text:{text_hash}")
            
            components.extend(slide_components)
        
        # Compute final hash
        version_string = "|".join(components)
        full_hash = hashlib.sha256(version_string.encode()).hexdigest()
        
        return full_hash[:16]
    
    # ========================================================================
    # PRIVATE HELPER METHODS
    # ========================================================================
    
    def _get_slide(self, index: int):
        """
        Get slide by index with validation.
        
        Args:
            index: Slide index (0-based)
            
        Returns:
            Slide object
            
        Raises:
            PowerPointAgentError: If no presentation loaded
            SlideNotFoundError: If index is out of range
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        slide_count = len(self.prs.slides)
        
        if not 0 <= index < slide_count:
            raise SlideNotFoundError(
                f"Slide index {index} out of range",
                details={"index": index, "slide_count": slide_count, "valid_range": f"0-{slide_count-1}"}
            )
        
        return self.prs.slides[index]
    
    def _get_shape(self, slide_index: int, shape_index: int):
        """
        Get shape by slide and shape index with validation.
        
        Args:
            slide_index: Slide index
            shape_index: Shape index on slide
            
        Returns:
            Shape object
            
        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
        """
        slide = self._get_slide(slide_index)
        
        shape_count = len(slide.shapes)
        
        if not 0 <= shape_index < shape_count:
            raise ShapeNotFoundError(
                f"Shape index {shape_index} out of range on slide {slide_index}",
                details={
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "shape_count": shape_count,
                    "valid_range": f"0-{shape_count-1}" if shape_count > 0 else "no shapes"
                }
            )
        
        return slide.shapes[shape_index]
    
    def _get_chart_shape(self, slide_index: int, chart_index: int):
        """
        Get chart shape by slide and chart index.
        
        Args:
            slide_index: Slide index
            chart_index: Chart index on slide (0-based among charts only)
            
        Returns:
            Chart shape object
            
        Raises:
            ChartNotFoundError: If chart not found
        """
        slide = self._get_slide(slide_index)
        
        chart_count = 0
        for shape in slide.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                if chart_count == chart_index:
                    return shape
                chart_count += 1
        
        raise ChartNotFoundError(
            f"Chart at index {chart_index} not found on slide {slide_index}",
            details={
                "slide_index": slide_index,
                "chart_index": chart_index,
                "charts_found": chart_count
            }
        )
    
    def _get_layout(self, layout_name: str):
        """
        Get layout by name with caching.
        
        Args:
            layout_name: Layout name
            
        Returns:
            Layout object
            
        Raises:
            LayoutNotFoundError: If layout doesn't exist
        """
        self._ensure_layout_cache()
        
        layout = self._layout_cache.get(layout_name)
        
        if layout is None:
            raise LayoutNotFoundError(
                f"Layout '{layout_name}' not found",
                details={"available_layouts": list(self._layout_cache.keys())}
            )
        
        return layout
    
    def _ensure_layout_cache(self) -> None:
        """Build layout cache if not already built."""
        if self._layout_cache is not None:
            return
        
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        
        self._layout_cache = {
            layout.name: layout
            for layout in self.prs.slide_layouts
        }
    
    def _get_placeholder_type_int(self, ph_type: Any) -> int:
        """Convert placeholder type to integer safely."""
        if ph_type is None:
            return 0
        if hasattr(ph_type, 'value'):
            return ph_type.value
        try:
            return int(ph_type)
        except (TypeError, ValueError):
            return 0
    
    def _copy_shape(self, source_shape, target_slide) -> None:
        """
        Copy shape to target slide.
        
        Args:
            source_shape: Shape to copy
            target_slide: Destination slide
        """
        # Handle pictures
        if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = source_shape.image.blob
                target_slide.shapes.add_picture(
                    BytesIO(blob),
                    source_shape.left, source_shape.top,
                    source_shape.width, source_shape.height
                )
            except Exception as e:
                logger.warning(f"Could not copy picture: {e}")
            return
        
        # Handle auto shapes and text boxes
        if source_shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX):
            try:
                # Get auto shape type, default to rectangle
                try:
                    auto_shape_type = source_shape.auto_shape_type
                except Exception:
                    auto_shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE
                
                new_shape = target_slide.shapes.add_shape(
                    auto_shape_type,
                    source_shape.left, source_shape.top,
                    source_shape.width, source_shape.height
                )
                
                # Copy text
                if source_shape.has_text_frame:
                    try:
                        new_shape.text_frame.text = source_shape.text_frame.text
                    except Exception:
                        pass
                
                # Copy fill
                try:
                    if source_shape.fill.type == 1:  # Solid fill
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                except Exception:
                    pass
                
            except Exception as e:
                logger.warning(f"Could not copy shape: {e}")
            return
        
        # Log unsupported shape types
        logger.debug(f"Shape type {source_shape.shape_type} not copied (not supported)")


# ============================================================================
# MODULE EXPORTS
# ============================================================================

__all__ = [
    # Main class
    "PowerPointAgent",
    
    # Exceptions
    "PowerPointAgentError",
    "SlideNotFoundError",
    "ShapeNotFoundError",
    "ChartNotFoundError",
    "LayoutNotFoundError",
    "ImageNotFoundError",
    "InvalidPositionError",
    "TemplateError",
    "ThemeError",
    "AccessibilityError",
    "AssetValidationError",
    "FileLockError",
    "PathValidationError",
    
    # Utility classes
    "FileLock",
    "PathValidator",
    "Position",
    "Size",
    "ColorHelper",
    "TemplateProfile",
    "AccessibilityChecker",
    "AssetValidator",
    
    # Enums
    "ShapeType",
    "ChartType",
    "TextAlignment",
    "VerticalAlignment",
    "BulletStyle",
    "ImageFormat",
    "ExportFormat",
    "ZOrderAction",
    "NotesMode",
    
    # Constants
    "SLIDE_WIDTH_INCHES",
    "SLIDE_HEIGHT_INCHES",
    "ANCHOR_POINTS",
    "CORPORATE_COLORS",
    "STANDARD_FONTS",
    "PLACEHOLDER_TYPE_NAMES",
    "TITLE_PLACEHOLDER_TYPES",
    "SUBTITLE_PLACEHOLDER_TYPE",
    "WCAG_CONTRAST_NORMAL",
    "WCAG_CONTRAST_LARGE",
    "EMU_PER_INCH",
    
    # Functions
    "get_placeholder_type_name",
    
    # Module metadata
    "__version__",
    "__author__",
    "__license__",
]
