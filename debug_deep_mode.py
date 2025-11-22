import sys
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from tools.ppt_capability_probe import PLACEHOLDER_TYPE_MAP, analyze_placeholder

def debug_deep_mode():
    print("Debugging Deep Mode...")
    prs = Presentation("test_probe.pptx")
    
    # Check Footer Type Code
    footer_type_code = None
    for type_code, type_name in PLACEHOLDER_TYPE_MAP.items():
        if type_name == 'FOOTER':
            footer_type_code = type_code
            print(f"FOOTER type code in MAP: {footer_type_code}")
            break
            
    # Instantiate first layout (Title Slide)
    layout = prs.slide_layouts[0]
    print(f"\nAnalyzing Layout: {layout.name}")
    
    # Standard Mode Check
    print("Standard Mode Placeholders:")
    for shape in layout.placeholders:
        print(f"  - {shape.name}: Type={shape.placeholder_format.type} ({PLACEHOLDER_TYPE_MAP.get(shape.placeholder_format.type)})")

    # Deep Mode Check
    print("\nDeep Mode (Instantiated) Placeholders:")
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            print(f"  - {shape.name}: Type={ph_type} ({PLACEHOLDER_TYPE_MAP.get(ph_type)})")
            
            if ph_type == footer_type_code:
                print("    -> MATCHES FOOTER CODE")
            else:
                print(f"    -> DOES NOT MATCH FOOTER CODE ({footer_type_code})")

if __name__ == "__main__":
    debug_deep_mode()
