#!/usr/bin/env python3
"""
PowerPoint Crop Image Tool
Crop an existing image on a slide

Usage:
    uv python ppt_crop_image.py --file deck.pptx --slide 0 --shape 1 --left 0.1 --right 0.1 --json
"""

import sys
import json
import argparse
from pathlib import Path
from typing import Dict, Any

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import (
    PowerPointAgent, PowerPointAgentError, SlideNotFoundError
)
from pptx.enum.shapes import MSO_SHAPE_TYPE

def crop_image(
    filepath: Path,
    slide_index: int,
    shape_index: int,
    left: float = 0.0,
    right: float = 0.0,
    top: float = 0.0,
    bottom: float = 0.0
) -> Dict[str, Any]:
    
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
        
    if not (0.0 <= left <= 1.0 and 0.0 <= right <= 1.0 and 0.0 <= top <= 1.0 and 0.0 <= bottom <= 1.0):
        raise ValueError("Crop values must be between 0.0 and 1.0")

    with PowerPointAgent(filepath) as agent:
        agent.open(filepath)
        
        total = agent.get_slide_count()
        if not 0 <= slide_index < total:
            raise SlideNotFoundError(f"Slide index {slide_index} out of range")
            
        slide = agent.prs.slides[slide_index]
        
        if not 0 <= shape_index < len(slide.shapes):
             raise ValueError(f"Shape index {shape_index} out of range")
             
        shape = slide.shapes[shape_index]
        
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape {shape_index} is not a picture")
            
        # Apply crop
        # python-pptx handles crop as a percentage of original size trimmed from edges
        if left > 0: shape.crop_left = left
        if right > 0: shape.crop_right = right
        if top > 0: shape.crop_top = top
        if bottom > 0: shape.crop_bottom = bottom
        
        agent.save()
        
    return {
        "status": "success",
        "file": str(filepath),
        "slide_index": slide_index,
        "shape_index": shape_index,
        "crop_applied": {
            "left": left,
            "right": right,
            "top": top,
            "bottom": bottom
        }
    }

def main():
    parser = argparse.ArgumentParser(description="Crop PowerPoint image")
    parser.add_argument('--file', required=True, type=Path, help='PowerPoint file path')
    parser.add_argument('--slide', required=True, type=int, help='Slide index (0-based)')
    parser.add_argument('--shape', required=True, type=int, help='Shape index (0-based)')
    parser.add_argument('--left', type=float, default=0.0, help='Crop from left (0.0-1.0)')
    parser.add_argument('--right', type=float, default=0.0, help='Crop from right (0.0-1.0)')
    parser.add_argument('--top', type=float, default=0.0, help='Crop from top (0.0-1.0)')
    parser.add_argument('--bottom', type=float, default=0.0, help='Crop from bottom (0.0-1.0)')
    parser.add_argument('--json', action='store_true', default=True, help='Output JSON')
    
    args = parser.parse_args()
    
    try:
        result = crop_image(
            filepath=args.file, 
            slide_index=args.slide, 
            shape_index=args.shape,
            left=args.left,
            right=args.right,
            top=args.top,
            bottom=args.bottom
        )
        print(json.dumps(result, indent=2))
        sys.exit(0)
    except Exception as e:
        print(json.dumps({"status": "error", "error": str(e), "error_type": type(e).__name__}, indent=2))
        sys.exit(1)

if __name__ == "__main__":
    main()
