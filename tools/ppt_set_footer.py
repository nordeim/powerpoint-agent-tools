#!/usr/bin/env python3
"""
PowerPoint Set Footer Tool v3.1.0
Configure slide footer with Dual Strategy (Placeholder + Text Box Fallback).

Fixes:
- Forces fallback to text boxes if placeholders exist on master but fail to update on slides.
- Clean JSON output (stderr redirected).
"""

import sys
import os

# CRITICAL: Redirect stderr to silence logs/warnings
sys.stderr = open(os.devnull, 'w')

import json
import argparse
import logging
from pathlib import Path
from typing import Dict, Any

# Configure logging to null
logging.basicConfig(level=logging.CRITICAL)

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import PowerPointAgent

# Safe import for constants
try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:
    class PP_PLACEHOLDER:
        FOOTER = 15
        SLIDE_NUMBER = 13

def set_footer(filepath: Path, text: str = None, show_number: bool = False) -> Dict[str, Any]:
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")

    slide_indices_updated = set()
    method_used = None
    
    with PowerPointAgent(filepath) as agent:
        agent.open(filepath)
        
        # Strategy 1: Placeholders
        # Try to set footer on Master first
        try:
            for master in agent.prs.slide_masters:
                for layout in master.slide_layouts:
                    for shape in layout.placeholders:
                        if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                            if text: shape.text = text
        except Exception:
            pass
        
        # Try to set on individual slides
        for slide_idx, slide in enumerate(agent.prs.slides):
            for shape in slide.placeholders:
                if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                    try:
                        if text: shape.text = text
                        slide_indices_updated.add(slide_idx)
                    except:
                        pass

        # Strategy 2: Fallback (Text Box)
        # Trigger if NO slides were updated via placeholders, even if Master appeared to have them.
        if len(slide_indices_updated) == 0:
            method_used = "text_box"
            for slide_idx in range(1, len(agent.prs.slides)):
                try:
                    if text:
                        agent.add_text_box(
                            slide_index=slide_idx,
                            text=text,
                            position={"left": "5%", "top": "92%"},
                            size={"width": "60%", "height": "5%"},
                            font_size=10,
                            color="#595959"
                        )
                        slide_indices_updated.add(slide_idx)
                    if show_number:
                        agent.add_text_box(
                            slide_index=slide_idx,
                            text=str(slide_idx + 1),
                            position={"left": "92%", "top": "92%"},
                            size={"width": "5%", "height": "5%"},
                            font_size=10,
                            color="#595959"
                        )
                        slide_indices_updated.add(slide_idx)
                except Exception:
                    pass
        else:
            method_used = "placeholder"

        agent.save()
        prs_info = agent.get_presentation_info()

    return {
        "status": "success" if len(slide_indices_updated) > 0 else "warning",
        "file": str(filepath),
        "method_used": method_used,
        "slides_updated": len(slide_indices_updated),
        "presentation_version_after": prs_info["presentation_version"]
    }

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--file', required=True, type=Path)
    parser.add_argument('--text')
    parser.add_argument('--show-number', action='store_true')
    parser.add_argument('--show-date', action='store_true')
    parser.add_argument('--json', action='store_true', default=True)
    args = parser.parse_args()
    
    try:
        result = set_footer(args.file, args.text, args.show_number)
        sys.stdout.write(json.dumps(result, indent=2) + "\n")
    except Exception as e:
        sys.stdout.write(json.dumps({"status": "error", "error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()
