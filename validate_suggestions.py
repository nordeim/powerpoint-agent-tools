import sys
import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

def test_enum_types():
    print("Testing PP_PLACEHOLDER types...")
    for name in dir(PP_PLACEHOLDER):
        if name.isupper():
            val = getattr(PP_PLACEHOLDER, name)
            print(f"{name}: type={type(val)}, value={val}")
            if hasattr(val, 'value'):
                print(f"  Has .value: {val.value}")

def test_theme_types():
    print("\nTesting Theme Types (using test_probe.pptx)...")
    if not os.path.exists("test_probe.pptx"):
        print("test_probe.pptx not found, creating one...")
        # Create a dummy one if needed, but we prefer one with a theme
        prs = Presentation()
        prs.save("test_probe.pptx")
    
    prs = Presentation("test_probe.pptx")
    slide_master = prs.slide_masters[0]
    
    # Font Scheme
    print("\nFont Scheme:")
    try:
        theme = slide_master.theme
        font_scheme = theme.font_scheme
        major = font_scheme.major_font
        print(f"major_font.latin type: {type(major.latin)}")
        print(f"major_font.latin value: {major.latin}")
        if hasattr(major.latin, 'typeface'):
            print(f"major_font.latin.typeface: {major.latin.typeface}")
    except Exception as e:
        print(f"Error checking fonts: {e}")

    # Color Scheme
    print("\nColor Scheme:")
    try:
        color_scheme = slide_master.theme.theme_color_scheme
        accent1 = color_scheme.accent1
        print(f"accent1 type: {type(accent1)}")
        print(f"accent1 value: {accent1}")
        if hasattr(accent1, 'rgb'):
            print(f"accent1.rgb: {accent1.rgb}")
    except Exception as e:
        print(f"Error checking colors: {e}")

if __name__ == "__main__":
    test_enum_types()
    test_theme_types()
